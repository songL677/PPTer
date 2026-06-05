from __future__ import annotations

from pathlib import Path

from parsers.base import DocumentPage, ParsedDocument, ParserError


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(
        paragraph.text.strip()
        for paragraph in shape.text_frame.paragraphs
        if paragraph.text and paragraph.text.strip()
    ).strip()


def parse_pptx(file_path: str | Path, file_name: str | None = None) -> ParsedDocument:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ParserError("缺少 python-pptx 依赖，请先运行：pip install -r requirements.txt") from exc

    path = Path(file_path)
    display_name = file_name or path.name

    try:
        presentation = Presentation(path)
    except Exception as exc:
        raise ParserError(f"PPTX 打开失败：{exc}") from exc

    pages: list[DocumentPage] = []
    warnings: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        title = ""
        title_shape = getattr(slide.shapes, "title", None)
        if title_shape is not None:
            title = _shape_text(title_shape)

        body_parts: list[str] = []
        for shape in slide.shapes:
            text = _shape_text(shape)
            if not text:
                continue
            if title and text == title:
                continue
            body_parts.append(text)

        slide_text = "\n".join(body_parts).strip()
        if not title and slide_text:
            title = slide_text.splitlines()[0][:60]
        if not title and not slide_text:
            warnings.append(f"第 {index} 张幻灯片没有提取到文字。")

        pages.append(
            DocumentPage(
                number=index,
                title=title or f"第 {index} 张幻灯片",
                text=slide_text,
                source_type="slide",
            )
        )

    if not pages:
        raise ParserError("PPTX 没有可读取的幻灯片。")

    if not any(page.text.strip() or page.title.strip() for page in pages):
        raise ParserError("PPTX 未提取到可用文字，可能主要由图片组成。")

    return ParsedDocument(
        file_name=display_name,
        file_type="PPTX",
        page_count=len(pages),
        status="解析完成",
        pages=pages,
        warnings=warnings,
    )
